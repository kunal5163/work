from pptx import Presentation

def convert_to_blank_layout(input_path, output_path):
    prs = Presentation(input_path)
    blank_layout = prs.slide_layouts[6]

    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # Remove default slide
    while len(new_prs.slides) > 0:
        r_id = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(r_id)
        del new_prs.slides._sldIdLst[0]

    for slide in prs.slides:
        new_slide = new_prs.slides.add_slide(blank_layout)
        for shape in slide.shapes:
            try:
                new_slide.shapes._spTree.insert_element_before(shape._element, 'p:extLst')
            except Exception as e:
                print(f"⚠️ Could not copy shape: {shape.name} — {e}")

    new_prs.save(output_path)
    print(f"✅ Slides converted to blank layout and saved to: {output_path}")

# Usage:
convert_to_blank_layout("input.pptx", "input_blank.pptx")
