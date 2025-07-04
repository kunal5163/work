import os
import json
import uuid
import base64
import zipfile
import subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_points(emu):
    return emu / 12700.0

def convert_pptx_to_pdf(pptx_path, output_pdf_path=None):
    if output_pdf_path is None:
        output_pdf_path = os.path.splitext(pptx_path)[0] + ".pdf"
    try:
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf", pptx_path,
            "--outdir", os.path.dirname(output_pdf_path) or "."
        ], check=True)
        print(f"✅ PDF generated at: {output_pdf_path}")
    except Exception as e:
        print(f"❌ PDF conversion failed: {e}")

def extract_combined_ppt_data(text_pptx_path, image_pptx_path, output_json_path="output_data.json", image_output_dir="extracted_images"):
    text_prs = Presentation(text_pptx_path)
    image_prs = Presentation(image_pptx_path)

    os.makedirs(image_output_dir, exist_ok=True)

    slides_json = []
    slide_width = text_prs.slide_width
    slide_height = text_prs.slide_height

    for slide_number, (text_slide, image_slide) in enumerate(zip(text_prs.slides, image_prs.slides)):
        slide_data = {
            "slide_number": slide_number + 1,
            "shapes": []
        }

        for shape_index, text_shape in enumerate(text_slide.shapes):
            shape_data = {
                "type": None,
                "name": text_shape.name,
                "position": {
                    "x_pt": emu_to_points(text_shape.left),
                    "y_pt": emu_to_points(text_shape.top)
                },
                "size": {
                    "width_pt": emu_to_points(text_shape.width),
                    "height_pt": emu_to_points(text_shape.height)
                }
            }

            if text_shape.has_text_frame:
                shape_data["type"] = "text"
                paragraphs_data = []
                full_text = ""
                for paragraph in text_shape.text_frame.paragraphs:
                    runs_data = []
                    para_text = ""
                    for run in paragraph.runs:
                        text = run.text or ""
                        para_text += text
                        full_text += text
                        run_data = {
                            "text": text,
                            "font_size_pt": emu_to_points(run.font.size) if run.font.size else None,
                            "font_name": run.font.name,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline
                        }
                        runs_data.append(run_data)

                    para_info = {
                        "alignment": str(paragraph.alignment).split('.')[-1].lower() if paragraph.alignment else None,
                        "runs": runs_data,
                        "line_spacing": emu_to_points(paragraph.line_spacing) if paragraph.line_spacing else None,
                        "space_before": emu_to_points(paragraph.space_before) if paragraph.space_before else None,
                        "space_after": emu_to_points(paragraph.space_after) if paragraph.space_after else None
                    }
                    paragraphs_data.append(para_info)

                shape_data["content"] = full_text.strip()
                shape_data["text_properties"] = {
                    "paragraphs": paragraphs_data,
                    "vertical_alignment": str(text_shape.text_frame.vertical_anchor).split('.')[-1].lower()
                        if text_shape.text_frame.vertical_anchor else None,
                    "margin_left_pt": emu_to_points(text_shape.text_frame.margin_left),
                    "margin_right_pt": emu_to_points(text_shape.text_frame.margin_right),
                    "margin_top_pt": emu_to_points(text_shape.text_frame.margin_top),
                    "margin_bottom_pt": emu_to_points(text_shape.text_frame.margin_bottom)
                }

            slide_data["shapes"].append(shape_data)

        for image_shape in image_slide.shapes:
            if image_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_data = {
                    "type": "image",
                    "name": image_shape.name,
                    "position": {
                        "x_pt": emu_to_points(image_shape.left),
                        "y_pt": emu_to_points(image_shape.top)
                    },
                    "size": {
                        "width_pt": emu_to_points(image_shape.width),
                        "height_pt": emu_to_points(image_shape.height)
                    }
                }

                image = image_shape.image
                image_ext = image.ext
                image_blob = image.blob

                if image_blob:
                    image_filename = f"slide{slide_number+1}_img_{uuid.uuid4().hex[:8]}.{image_ext}"
                    image_path = os.path.join(image_output_dir, image_filename)
                    try:
                        with open(image_path, 'wb') as f:
                            f.write(image_blob)
                        with open(image_path, 'rb') as img_file:
                            encoded_string = base64.b64encode(img_file.read()).decode('utf-8')

                        shape_data["image_metadata"] = {
                            "content_type": image.content_type,
                            "ext": image_ext,
                            "filename": image_filename,
                            "saved_path": os.path.abspath(image_path),
                            "thumbnail_base64": encoded_string
                        }
                    except Exception as e:
                        shape_data["image_metadata"] = {"error": str(e)}

                    slide_data["shapes"].append(shape_data)

        slides_json.append(slide_data)

    presentation_data = {
        "slide_width_emu": slide_width,
        "slide_height_emu": slide_height,
        "slides": slides_json
    }

    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, indent=4, ensure_ascii=False)
    print(f"\n✅ JSON saved to: {output_json_path}")

    image_files = [f for f in os.listdir(image_output_dir) if f.lower().endswith(('.jpg', '.jpeg', '.png', '.bmp', '.gif'))]
    zip_path = f"{image_output_dir}.zip"
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for file in image_files:
            zipf.write(os.path.join(image_output_dir, file), file)
    print(f"✅ Zipped images saved to: {zip_path}")

# Example usage
extract_combined_ppt_data("input_blank.pptx", "input.pptx")
