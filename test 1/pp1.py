import os
import json
import uuid
import base64
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_points(emu):
    return emu / 12700.0

def extract_images_from_pptx(pptx_path, image_output_dir):
    prs = Presentation(pptx_path)
    os.makedirs(image_output_dir, exist_ok=True)
    images = []
    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_position = {
                    "x_pt": emu_to_points(shape.left),
                    "y_pt": emu_to_points(shape.top)
                }
                shape_size = {
                    "width_pt": emu_to_points(shape.width),
                    "height_pt": emu_to_points(shape.height)
                }
                image = shape.image
                image_ext = image.ext
                image_blob = image.blob
                image_filename = f"slide{slide_index+1}_img_{uuid.uuid4().hex[:8]}.{image_ext}"
                image_path = os.path.join(image_output_dir, image_filename)
                with open(image_path, 'wb') as f:
                    f.write(image_blob)
                with open(image_path, 'rb') as img_file:
                    encoded_string = base64.b64encode(img_file.read()).decode('utf-8')
                images.append({
                    "slide_index": slide_index,
                    "name": shape.name,
                    "position": shape_position,
                    "size": shape_size,
                    "image_metadata": {
                        "content_type": image.content_type,
                        "ext": image_ext,
                        "filename": image_filename,
                        "saved_path": os.path.abspath(image_path),
                        "thumbnail_base64": encoded_string
                    }
                })
    return images

def combine_blank_with_images(blank_json_path, pptx_path, output_json="output_data.json", image_output_dir="extracted_images"):
    with open(blank_json_path, 'r', encoding='utf-8') as f:
        blank_data = json.load(f)
    images = extract_images_from_pptx(pptx_path, image_output_dir)
    # Remove all image shapes from blank_data
    for slide in blank_data["slides"]:
        slide["shapes"] = [s for s in slide["shapes"] if s["type"] != "image"]
    # Add images to the correct slide
    for img in images:
        slide = blank_data["slides"][img["slide_index"]]
        image_shape = {
            "type": "image",
            "name": img["name"],
            "position": img["position"],
            "size": img["size"],
            "image_metadata": img["image_metadata"]
        }
        slide["shapes"].append(image_shape)
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(blank_data, f, indent=4, ensure_ascii=False)
    print(f"\n✅ Combined JSON with images saved to: {output_json}")
    # Zip image folder
    zip_path = f"{image_output_dir}.zip"
    image_files = [f for f in os.listdir(image_output_dir) if f.lower().endswith((".jpg", ".jpeg", ".png", ".bmp", ".gif"))]
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for file in image_files:
            zipf.write(os.path.join(image_output_dir, file), file)
    print(f"✅ Zipped images saved to: {zip_path}")

# Run
if __name__ == "__main__":
    combine_blank_with_images("blank_structure.json", "input.pptx")
