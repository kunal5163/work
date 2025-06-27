import json
import os
import zipfile
import io
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

def points_to_emu(pt):
    return round(pt * 12700)

def extract_image_from_zip(zip_path, image_filename):
    with zipfile.ZipFile(zip_path, 'r') as zipf:
        with zipf.open(image_filename) as img_file:
            return io.BytesIO(img_file.read())

def create_ppt_from_json_with_zip(json_path, zip_path, output_pptx="rebuilt_from_layout.pptx"):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()

    if "slide_width_emu" in data and "slide_height_emu" in data:
        prs.slide_width = data["slide_width_emu"]
        prs.slide_height = data["slide_height_emu"]

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY
    }

    for slide_info in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for shape in slide_info.get("shapes", []):
            x = points_to_emu(shape["position"]["x_pt"])
            y = points_to_emu(shape["position"]["y_pt"])
            width = points_to_emu(shape["size"]["width_pt"])
            height = points_to_emu(shape["size"]["height_pt"])

            if shape["type"] == "text":
                textbox = slide.shapes.add_textbox(left=x, top=y, width=width, height=height)
                text_frame = textbox.text_frame
                text_frame.clear()

                props = shape.get("text_properties", {})
                if props.get("vertical_alignment"):
                    try:
                        text_frame.vertical_anchor = getattr(MSO_VERTICAL_ANCHOR, props["vertical_alignment"].upper())
                    except:
                        pass

                text_frame.margin_left = points_to_emu(props.get("margin_left_pt", 0))
                text_frame.margin_right = points_to_emu(props.get("margin_right_pt", 0))
                text_frame.margin_top = points_to_emu(props.get("margin_top_pt", 0))
                text_frame.margin_bottom = points_to_emu(props.get("margin_bottom_pt", 0))

                rendered_lines = shape.get("rendered_lines")
                if rendered_lines:
                    for i, line in enumerate(rendered_lines):
                        para = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                        run = para.add_run()
                        run.text = line
                else:
                    paragraphs = props.get("paragraphs", [])
                    for para_index, para_data in enumerate(paragraphs):
                        para = text_frame.paragraphs[0] if para_index == 0 else text_frame.add_paragraph()

                        if para_data.get("alignment"):
                            para.alignment = align_map.get(para_data["alignment"], PP_ALIGN.LEFT)
                        if para_data.get("line_spacing"):
                            para.line_spacing = Pt(para_data["line_spacing"])
                        if para_data.get("space_before"):
                            para.space_before = Pt(para_data["space_before"])
                        if para_data.get("space_after"):
                            para.space_after = Pt(para_data["space_after"])

                        for run_data in para_data.get("runs", []):
                            raw_text = run_data.get("text", "")
                            run = para.add_run()
                            run.text = raw_text
                            if run_data.get("font_size_pt"):
                                run.font.size = Pt(run_data["font_size_pt"])
                            if run_data.get("font_name"):
                                run.font.name = run_data["font_name"]
                            if run_data.get("bold") is not None:
                                run.font.bold = run_data["bold"]
                            if run_data.get("italic") is not None:
                                run.font.italic = run_data["italic"]
                            if run_data.get("underline") is not None:
                                run.font.underline = run_data["underline"]

            elif shape["type"] == "image":
                metadata = shape.get("image_metadata", {})
                filename = metadata.get("filename")
                if filename:
                    try:
                        image_stream = extract_image_from_zip(zip_path, filename)
                        slide.shapes.add_picture(image_stream, x, y, width=width, height=height)
                        print(f"[✓] Added image: {filename}")
                    except Exception as e:
                        print(f"[✗] Could not add image {filename}: {e}")

    prs.save(output_pptx)
    print(f"\n✅ Presentation saved to: {output_pptx}")

# To run:
# create_ppt_from_json_with_zip("output_with_layout.json", "extracted_images.zip")      
# Example usage
if __name__ == "__main__":  
    json_file = "output_data.json"   # your JSON file
    zip_file = "extracted_images.zip" # your ZIP file with images
    output_pptx_file = "rebuilt_presentation.pptx"  # output PPTX file

    create_ppt_from_json_with_zip(json_file, zip_file, output_pptx_file)