import os
import json
import zipfile
from pptx import Presentation
from pptx.util import Pt

def points_to_emu(pt):
    return int(pt * 12700)

def extract_images_from_zip(zip_path, extract_dir):
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)

def get_text_from_shape(shape):
    # Try to get text from 'content', else reconstruct from text_properties
    if shape.get("content"):
        return shape["content"]
    # Try to reconstruct from text_properties
    text_props = shape.get("text_properties")
    if text_props and "paragraphs" in text_props:
        paras = []
        for para in text_props["paragraphs"]:
            para_text = "".join(run.get("text", "") for run in para.get("runs", []))
            paras.append(para_text)
        return "\n".join(paras)
    return ""

def build_pptx_from_json(json_path, images_dir, output_pptx="rebuilt_presentation.pptx"):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = points_to_emu(data["slide_width_emu"] / 12700.0)
    prs.slide_height = points_to_emu(data["slide_height_emu"] / 12700.0)

    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
        for shape in slide_data["shapes"]:
            if shape["type"] == "text":
                left = points_to_emu(shape["position"]["x_pt"])
                top = points_to_emu(shape["position"]["y_pt"])
                width = points_to_emu(shape["size"]["width_pt"])
                height = points_to_emu(shape["size"]["height_pt"])
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = get_text_from_shape(shape)
            elif shape["type"] == "image":
                img_path = os.path.join(images_dir, shape["image_metadata"]["filename"])
                left = points_to_emu(shape["position"]["x_pt"])
                top = points_to_emu(shape["position"]["y_pt"])
                width = points_to_emu(shape["size"]["width_pt"])
                height = points_to_emu(shape["size"]["height_pt"])
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, left, top, width, height)
                else:
                    print(f"Image not found: {img_path}")

    prs.save(output_pptx)
    print(f"✅ PPTX generated: {output_pptx}")

if __name__ == "__main__":
    json_path = "output_data.json"
    zip_path = "extracted_images.zip"
    images_dir = "extracted_images"
    extract_images_from_zip(zip_path, images_dir)
    build_pptx_from_json(json_path, images_dir)
