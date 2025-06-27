import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_points(emu):
    return emu / 12700.0

def extract_shapes_to_json(input_path, json_output="blank_structure.json"):
    prs = Presentation(input_path)
    blank_layout = prs.slide_layouts[6]

    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    while len(new_prs.slides) > 0:
        r_id = new_prs.slides._sldIdLst[0].rId
        new_prs.part.drop_rel(r_id)
        del new_prs.slides._sldIdLst[0]

    slides_json = []
    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": slide_num + 1,
            "shapes": []
        }
        new_slide = new_prs.slides.add_slide(blank_layout)

        for shape in slide.shapes:
            shape_data = {
                "name": shape.name,
                "position": {
                    "x_pt": emu_to_points(shape.left),
                    "y_pt": emu_to_points(shape.top)
                },
                "size": {
                    "width_pt": emu_to_points(shape.width),
                    "height_pt": emu_to_points(shape.height)
                }
            }

            if shape.has_text_frame:
                shape_data["type"] = "text"
                # Extract actual text content
                text_content = shape.text_frame.text if shape.text_frame else ""
                shape_data["content"] = text_content
                textbox = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                tf = textbox.text_frame
                tf.clear()
                tf.paragraphs[0].add_run().text = "[Text Placeholder]"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_data["type"] = "image"
                placeholder = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
                placeholder.text_frame.text = "[Image Placeholder]"
            else:
                continue

            slide_data["shapes"].append(shape_data)

        slides_json.append(slide_data)

    presentation_data = {
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slides": slides_json
    }

    with open(json_output, 'w', encoding='utf-8') as f:
        json.dump(presentation_data, f, indent=4, ensure_ascii=False)
    print(f"✅ JSON with shape layout saved to: {json_output}")

    new_prs.save("input_blank.pptx")
    print(f"✅ Blank layout presentation saved to: input_blank.pptx")

# Run
extract_shapes_to_json("input.pptx")
