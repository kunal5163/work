import os
import json
import uuid
import base64
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def emu_to_points(emu):
    return emu / 12700.0

def is_match(pos1, size1, pos2, size2, tolerance=1.5):
    """Check if two shapes have roughly the same position and size (in points)."""
    for key in pos1:
        if abs(pos1[key] - pos2[key]) > tolerance:
            return False
    for key in size1:
        if abs(size1[key] - size2[key]) > tolerance:
            return False
    return True

def update_blank_json_with_images_precise(blank_json_path, original_pptx_path, output_json="output_data.json", image_output_dir="extracted_images"):
    with open(blank_json_path, 'r', encoding='utf-8') as f:
        json_data = json.load(f)

    prs = Presentation(original_pptx_path)
    os.makedirs(image_output_dir, exist_ok=True)

    for slide_index, slide in enumerate(prs.slides):
        if slide_index >= len(json_data["slides"]):
            continue

        slide_data = json_data["slides"][slide_index]

        # Remove image placeholders
        slide_data["shapes"] = [s for s in slide_data["shapes"] if s["type"] != "image"]

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_position = {
                    "x_pt": emu_to_points(shape.left),
                    "y_pt": emu_to_points(shape.top)
                }
                shape_size = {
                    "width_pt": emu_to_points(shape.width),
                    "height_pt": emu_to_points(shape.height)
                }

                image = shape.image
                image_ext = image.ext
                image_blob = image.blob

                if image_blob:
                    image_filename = f"slide{slide_index+1}_img_{uuid.uuid4().hex[:8]}.{image_ext}"
                    image_path = os.path.join(image_output_dir, image_filename)

                    try:
                        with open(image_path, 'wb') as f:
                            f.write(image_blob)

                        with open(image_path, 'rb') as img_file:
                            encoded_string = base64.b64encode(img_file.read()).decode('utf-8')

                        image_shape_data = {
                            "type": "image",
                            "name": shape.name,
                            "position": shape_position,
                            "size": shape_size,
                            "image_metadata": {
                                "content_type": image.content_type,
                                "ext": image_ext,
                                "filename": image_filename,
                                "saved_path": os.path.abspath(image_path),
                                "thumbnail_base64": encoded_string
                            }
                        }

                        # Match to JSON by position/size if desired (optional)
                        matched = False
                        for json_shape in slide_data["shapes"]:
                            if json_shape["type"] == "image":
                                if is_match(json_shape["position"], json_shape["size"], shape_position, shape_size):
                                    matched = True
                                    break

                        slide_data["shapes"].append(image_shape_data)

                    except Exception as e:
                        print(f"⚠️ Could not extract image: {e}")

    # Save updated JSON
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(json_data, f, indent=4, ensure_ascii=False)
    print(f"\n✅ Final JSON with precise image metadata saved to: {output_json}")

    # Zip image folder
    zip_path = f"{image_output_dir}.zip"
    image_files = [f for f in os.listdir(image_output_dir) if f.lower().endswith(('.jpg', '.jpeg', '.png', '.bmp', '.gif'))]
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for file in image_files:
            zipf.write(os.path.join(image_output_dir, file), file)
    print(f"✅ Zipped images saved to: {zip_path}")

# Run
if __name__ == "__main__":
    update_blank_json_with_images_precise("blank_structure.json", "input.pptx")
